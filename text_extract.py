"""Best-effort local text extraction for common evidence file types."""
from __future__ import annotations

import csv
import html
import json
import re
import zipfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Iterable

from common import EXTRACT_MAX_CHARS

TEXTLIKE_EXTS = {".txt", ".md", ".markdown", ".log", ".json", ".xml", ".html", ".htm", ".py", ".js", ".css", ".sql"}


class ExtractionResult(dict):
    """Simple dict subclass for readability."""


def _cap(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + "\n\n[TRUNCATED: extracted text exceeded configured character limit.]"


def _append_limited(parts: list[str], text: str, max_chars: int) -> bool:
    """Append text and return False once the aggregate is at/over limit."""
    if not text:
        return True
    current = sum(len(p) for p in parts)
    remaining = max_chars - current
    if remaining <= 0:
        return False
    if len(text) > remaining:
        parts.append(text[:remaining])
        return False
    parts.append(text)
    return True


def _read_text_file(path: Path, max_chars: int) -> str:
    # Read a bit beyond the cap so decoded multibyte characters do not get split too often.
    with path.open("rb") as f:
        data = f.read(max_chars * 2)
    return _cap(data.decode("utf-8", errors="replace"), max_chars)


def _extract_pdf(path: Path, max_chars: int) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        keep_going = _append_limited(parts, f"\n\n--- PDF page {idx} ---\n{text}", max_chars)
        if not keep_going:
            break
    return _cap("".join(parts).strip(), max_chars)


def _xml_payload_to_text(data: bytes) -> str:
    raw = data.decode("utf-8", errors="replace")
    # Preserve paragraph/table boundaries from Office XML reasonably well.
    raw = re.sub(r"</w:p>|</a:p>|</w:tr>|</a:tr>", "\n", raw)
    raw = re.sub(r"</w:tc>|</a:tc>", "\t", raw)
    raw = re.sub(r"<[^>]+>", " ", raw)
    raw = html.unescape(raw)
    raw = re.sub(r"[ \t]+", " ", raw)
    raw = re.sub(r"\n\s+", "\n", raw)
    raw = re.sub(r"\n{3,}", "\n\n", raw)
    return raw.strip()


def _extract_docx_zip(path: Path, max_chars: int) -> str:
    """Fallback DOCX extractor that reads the zipped XML directly."""
    parts: list[str] = []
    wanted = (
        "word/document.xml",
        "word/header",
        "word/footer",
        "word/footnotes",
        "word/endnotes",
        "word/comments",
    )
    with zipfile.ZipFile(path) as zf:
        for name in zf.namelist():
            if not (name == wanted[0] or name.startswith(wanted[1:])):
                continue
            try:
                text = _xml_payload_to_text(zf.read(name))
            except Exception:
                continue
            if text:
                if not _append_limited(parts, f"\n\n--- {name} ---\n{text}", max_chars):
                    break
    return _cap("".join(parts).strip(), max_chars)


def _extract_docx(path: Path, max_chars: int) -> str:
    try:
        import docx

        doc = docx.Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            if not _append_limited(parts, para.text + "\n", max_chars):
                break
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if not _append_limited(parts, " | ".join(cells) + "\n", max_chars):
                    return _cap("".join(parts), max_chars)
        text = _cap("".join(parts), max_chars)
        if text.strip():
            return text
    except Exception:
        # Fall through to the zipped-XML fallback below.
        pass

    fallback = _extract_docx_zip(path, max_chars)
    if fallback.strip():
        return fallback
    # Let the caller record an extraction error if neither approach worked.
    raise ValueError("No text could be extracted from DOCX with python-docx or XML fallback.")


def _extract_xlsx(path: Path, max_chars: int) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            if not _append_limited(parts, f"\n\n--- Sheet: {ws.title} ---\n", max_chars):
                break
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v) for v in row]
                # Skip fully empty rows.
                if not any(values):
                    continue
                if not _append_limited(parts, "\t".join(values) + "\n", max_chars):
                    return _cap("".join(parts), max_chars)
    finally:
        wb.close()
    return _cap("".join(parts), max_chars)


def _extract_csv(path: Path, max_chars: int) -> str:
    parts: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        # Sniffing huge or malformed CSVs can be expensive; treat as rows and join.
        reader = csv.reader(f)
        for row in reader:
            if not _append_limited(parts, "\t".join(row) + "\n", max_chars):
                break
    return _cap("".join(parts), max_chars)


def _extract_pptx(path: Path, max_chars: int) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        if not _append_limited(parts, f"\n\n--- Slide {idx} ---\n", max_chars):
            break
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                if not _append_limited(parts, text + "\n", max_chars):
                    return _cap("".join(parts), max_chars)
    return _cap("".join(parts), max_chars)


def _walk_email_parts(msg) -> Iterable[str]:
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            disp = (part.get_content_disposition() or "").lower()
            if disp == "attachment":
                continue
            if content_type == "text/plain":
                try:
                    yield part.get_content()
                except Exception:
                    payload = part.get_payload(decode=True) or b""
                    yield payload.decode(part.get_content_charset() or "utf-8", errors="replace")
    else:
        try:
            yield msg.get_content()
        except Exception:
            payload = msg.get_payload(decode=True) or b""
            yield payload.decode(msg.get_content_charset() or "utf-8", errors="replace")


def _extract_eml(path: Path, max_chars: int) -> str:
    with path.open("rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    headers = {
        "Subject": msg.get("subject", ""),
        "From": msg.get("from", ""),
        "To": msg.get("to", ""),
        "Cc": msg.get("cc", ""),
        "Date": msg.get("date", ""),
    }
    parts: list[str] = [json.dumps(headers, indent=2), "\n\n--- Body ---\n"]
    for body in _walk_email_parts(msg):
        if not _append_limited(parts, str(body) + "\n", max_chars):
            break
    return _cap("".join(parts), max_chars)


def extract_text(path: str | Path, max_chars: int = EXTRACT_MAX_CHARS) -> ExtractionResult:
    """Return extracted text plus status metadata.

    Unsupported binary formats return a useful status and empty text so Ollama can still
    receive file metadata. This keeps the categorization job from failing on videos,
    audio, images, legacy Office files, and archives.
    """
    p = Path(path)
    ext = p.suffix.lower()
    try:
        if ext in TEXTLIKE_EXTS:
            text = _read_text_file(p, max_chars)
            status = "ok"
        elif ext in {".csv", ".tsv"}:
            text = _extract_csv(p, max_chars)
            status = "ok"
        elif ext == ".pdf":
            text = _extract_pdf(p, max_chars)
            status = "ok"
        elif ext == ".docx":
            text = _extract_docx(p, max_chars)
            status = "ok"
        elif ext in {".xlsx", ".xlsm"}:
            text = _extract_xlsx(p, max_chars)
            status = "ok"
        elif ext == ".pptx":
            text = _extract_pptx(p, max_chars)
            status = "ok"
        elif ext == ".eml":
            text = _extract_eml(p, max_chars)
            status = "ok"
        else:
            text = ""
            status = "unsupported_metadata_only"
        return ExtractionResult(status=status, text=text, chars=len(text), error="")
    except Exception as e:  # Return metadata-only rather than stopping the batch.
        return ExtractionResult(
            status="extract_error_metadata_only",
            text="",
            chars=0,
            error=f"{type(e).__name__}: {e}",
        )
